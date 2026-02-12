"""
Ingestor for PDF and PPTX lecture slides.
"""

from pathlib import Path
from typing import List
import fitz  # PyMuPDF
from pptx import Presentation

from src.core.indexing.ingestors.base import BaseIngestor, DocumentChunk


class SlidesIngestor(BaseIngestor):
    """Ingest PDF and PPTX lecture slides."""
    
    def can_ingest(self, path: Path) -> bool:
        """Check if file is PDF or PPTX."""
        suffix = path.suffix.lower()
        return suffix in ['.pdf', '.pptx', '.ppt']
    
    def ingest(self, path: Path) -> List[DocumentChunk]:
        """Ingest slides - one chunk per slide."""
        suffix = path.suffix.lower()
        
        if suffix == '.pdf':
            return self._ingest_pdf(path)
        elif suffix in ['.pptx', '.ppt']:
            return self._ingest_pptx(path)
        else:
            raise ValueError(f"Unsupported slide format: {suffix}")
    
    def _ingest_pdf(self, path: Path) -> List[DocumentChunk]:
        """Ingest PDF slides."""
        chunks = []
        
        try:
            doc = fitz.open(path)
            source_name = path.stem
            
            for page_num, page in enumerate(doc):
                # Extract text from page
                text = page.get_text()
                
                # Extract page title (first line or heading)
                lines = text.split('\n')
                title = lines[0].strip() if lines else f"Slide {page_num + 1}"
                
                # Extract speaker notes if present
                notes = ""
                # PyMuPDF doesn't directly extract notes, but we can try annotations
                # For now, just use main text
                
                # Combine title and content
                slide_text = f"{title}\n\n{text}".strip()
                
                if slide_text:
                    chunk = self._create_chunk(
                        text=slide_text,
                        source_type="lecture_slide",
                        source_name=source_name,
                        page=page_num + 1,
                        slide_title=title,
                        total_pages=len(doc)
                    )
                    chunks.append(chunk)
            
            doc.close()
        
        except Exception as e:
            raise ValueError(f"Failed to ingest PDF {path}: {str(e)}") from e
        
        return chunks
    
    def _ingest_pptx(self, path: Path) -> List[DocumentChunk]:
        """Ingest PPTX slides."""
        chunks = []
        
        try:
            prs = Presentation(path)
            source_name = path.stem
            
            for slide_num, slide in enumerate(prs.slides):
                # Extract text from all shapes
                text_parts = []
                title = ""
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text = shape.text.strip()
                        if text:
                            # First text box is usually title
                            if not title and len(text) < 100:
                                title = text
                            else:
                                text_parts.append(text)
                
                # Combine title and content
                slide_text = f"{title}\n\n" + "\n\n".join(text_parts) if title else "\n\n".join(text_parts)
                
                # Extract speaker notes
                notes = ""
                if slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes = notes_slide.notes_text_frame.text
                
                if notes:
                    slide_text += f"\n\n[Notes: {notes}]"
                
                if slide_text.strip():
                    chunk = self._create_chunk(
                        text=slide_text.strip(),
                        source_type="lecture_slide",
                        source_name=source_name,
                        page=slide_num + 1,
                        slide_title=title or f"Slide {slide_num + 1}",
                        total_pages=len(prs.slides),
                        has_notes=bool(notes)
                    )
                    chunks.append(chunk)
        
        except Exception as e:
            raise ValueError(f"Failed to ingest PPTX {path}: {str(e)}") from e
        
        return chunks
